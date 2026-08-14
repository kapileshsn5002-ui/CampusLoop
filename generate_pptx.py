import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    NAVY = RGBColor(10, 37, 64)       # #0A2540
    DARK_BLUE = RGBColor(15, 76, 129)  # #0F4C81
    ACCENT_BLUE = RGBColor(37, 99, 235)# #2563EB
    LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC
    CARD_BG = RGBColor(255, 255, 255) # #FFFFFF
    CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0
    TEXT_DARK = RGBColor(30, 41, 59)  # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
    SUCCESS_GREEN = RGBColor(16, 185, 129) # #10B981
    DANGER_RED = RGBColor(239, 68, 68) # #EF4444
    PILL_BG = RGBColor(239, 246, 255) # #EFF6FF

    def apply_background(slide, color=LIGHT_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CAMPUSLOOP PRESENTATION", num=1):
        # Header category badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_BLUE

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.65), Inches(10.0), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY

        # Top Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

        # Footer Slide Number
        footer_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.3))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.alignment = PP_ALIGN.RIGHT
        p_foot.text = f"{num} / 10"
        p_foot.font.size = Pt(11)
        p_foot.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 1 — TITLE
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1, LIGHT_BG)

    # Decorative Left Banner
    left_banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), Inches(7.5))
    left_banner.fill.solid()
    left_banner.fill.fore_color.rgb = ACCENT_BLUE
    left_banner.line.fill.background()

    # Academic Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(4.2), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PILL_BG
    badge.line.color.rgb = ACCENT_BLUE
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    p_b.text = "ACADEMIC PROJECT PRESENTATION (2025–2026)"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_BLUE

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(7.5), Inches(1.2))
    tf_t = t_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = "CampusLoop"
    p_t.font.size = Pt(44)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY

    # Subtitle
    st_box = s1.shapes.add_textbox(Inches(0.75), Inches(2.4), Inches(7.5), Inches(0.6))
    tf_st = st_box.text_frame
    p_st = tf_st.paragraphs[0]
    p_st.text = "Employee Leave Management System"
    p_st.font.size = Pt(22)
    p_st.font.bold = True
    p_st.font.color.rgb = DARK_BLUE

    # Tagline
    tag_box = s1.shapes.add_textbox(Inches(0.75), Inches(3.0), Inches(7.5), Inches(0.5))
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "Web-Based Employee Leave Management & Approval Platform"
    p_tag.font.size = Pt(14)
    p_tag.font.italic = True
    p_tag.font.color.rgb = TEXT_MUTED

    # Tech Stack Box
    tech_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.7), Inches(7.2), Inches(0.5))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = CARD_BG
    tech_box.line.color.rgb = CARD_BORDER
    tf_tech = tech_box.text_frame
    p_tech = tf_tech.paragraphs[0]
    p_tech.alignment = PP_ALIGN.CENTER
    p_tech.text = "React.js  |  Spring Boot  |  MySQL  |  REST API  |  Playwright"
    p_tech.font.size = Pt(12)
    p_tech.font.bold = True
    p_tech.font.color.rgb = ACCENT_BLUE

    # Right Card: Student Details Placeholder
    card_info = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(0.8), Inches(4.1), Inches(5.8))
    card_info.fill.solid()
    card_info.fill.fore_color.rgb = CARD_BG
    card_info.line.color.rgb = CARD_BORDER
    
    tf_ci = card_info.text_frame
    tf_ci.word_wrap = True
    tf_ci.margin_left = Inches(0.3)
    tf_ci.margin_top = Inches(0.4)
    
    p_hdr = tf_ci.paragraphs[0]
    p_hdr.text = "STUDENT DETAILS"
    p_hdr.font.size = Pt(14)
    p_hdr.font.bold = True
    p_hdr.font.color.rgb = NAVY
    
    details = [
        ("Student Name:", "[Student Name]"),
        ("Register Number:", "[Register Number]"),
        ("Department:", "[Department]"),
        ("College Name:", "[College Name]"),
        ("Academic Year:", "2025–2026")
    ]
    for lbl, val in details:
        p1 = tf_ci.add_paragraph()
        p1.space_before = Pt(14)
        p1.text = lbl
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED
        
        p2 = tf_ci.add_paragraph()
        p2.text = val
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_DARK

    # Slide 1 Footer
    f1 = s1.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(5.0), Inches(0.3))
    f1.text_frame.paragraphs[0].text = "Slide 1 / 10"
    f1.text_frame.paragraphs[0].font.size = Pt(11)
    f1.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2 — PROBLEM STATEMENT
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Problem Statement", num=2)

    # Left Card: Traditional Limitations
    c2_left = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.1))
    c2_left.fill.solid()
    c2_left.fill.fore_color.rgb = CARD_BG
    c2_left.line.color.rgb = CARD_BORDER
    tf_l2 = c2_left.text_frame
    tf_l2.margin_left = Inches(0.3)
    tf_l2.margin_top = Inches(0.3)
    
    p = tf_l2.paragraphs[0]
    p.text = "TRADITIONAL LEAVE MANAGEMENT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DANGER_RED

    points_l2 = [
        "Manual paper-based leave applications",
        "Informal & unorganized approval workflows",
        "Difficulty tracking real-time leave balances",
        "Lack of centralized employee record-keeping",
        "Delayed manager responses & approval bottlenecks",
        "Limited visibility of application status for employees",
        "Increased risk of data-entry & accounting errors"
    ]
    for pt in points_l2:
        p = tf_l2.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"•  {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK

    # Right Card: Core Problem Definition
    c2_right = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.1))
    c2_right.fill.solid()
    c2_right.fill.fore_color.rgb = CARD_BG
    c2_right.line.color.rgb = CARD_BORDER
    tf_r2 = c2_right.text_frame
    tf_r2.margin_left = Inches(0.4)
    tf_r2.margin_right = Inches(0.4)
    tf_r2.margin_top = Inches(0.4)
    
    p = tf_r2.paragraphs[0]
    p.text = "CORE PROBLEM DEFINITION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p_prob = tf_r2.add_paragraph()
    p_prob.space_before = Pt(14)
    p_prob.text = "Employees and managers require a centralized digital web platform where leave requests can be submitted, tracked, reviewed, and approved efficiently without administrative friction."
    p_prob.font.size = Pt(14)
    p_prob.font.color.rgb = TEXT_DARK
    p_prob.font.bold = True

    # Traditional Flow Box at bottom
    flow2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.9))
    flow2.fill.solid()
    flow2.fill.fore_color.rgb = RGBColor(254, 242, 242)
    flow2.line.color.rgb = RGBColor(252, 165, 165)
    tf_fl2 = flow2.text_frame
    p_fl2 = tf_fl2.paragraphs[0]
    p_fl2.alignment = PP_ALIGN.CENTER
    p_fl2.text = "Traditional Flow (Flawed):  Employee  ➔  Manual Paper Form  ➔  Manager Desk  ➔  Delayed Approval  ➔  Manual Records"
    p_fl2.font.size = Pt(12)
    p_fl2.font.bold = True
    p_fl2.font.color.rgb = DANGER_RED

    # -------------------------------------------------------------
    # SLIDE 3 — PROPOSED SOLUTION
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "Proposed Solution", num=3)

    # Intro Header Card
    c3_top = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    c3_top.fill.solid()
    c3_top.fill.fore_color.rgb = PILL_BG
    c3_top.line.color.rgb = ACCENT_BLUE
    tf_top3 = c3_top.text_frame
    p3_top = tf_top3.paragraphs[0]
    p3_top.alignment = PP_ALIGN.CENTER
    p3_top.text = "CampusLoop is a web-based employee leave management system designed to digitize, automate, and simplify the complete organizational leave workflow."
    p3_top.font.size = Pt(13)
    p3_top.font.bold = True
    p3_top.font.color.rgb = DARK_BLUE

    # Left Grid: Key Solution Points
    c3_left = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.5), Inches(6.5), Inches(4.3))
    c3_left.fill.solid()
    c3_left.fill.fore_color.rgb = CARD_BG
    c3_left.line.color.rgb = CARD_BORDER
    tf_l3 = c3_left.text_frame
    tf_l3.margin_left = Inches(0.3)
    tf_l3.margin_top = Inches(0.3)
    
    p = tf_l3.paragraphs[0]
    p.text = "KEY SOLUTION ADVANTAGES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sol_points = [
        "Secure Employee & Manager Authentication",
        "Interactive Employee Dashboard with Leave History",
        "Online Application Submission with Instant Validation",
        "Real-Time Leave Balance Calculation & Tracking",
        "Streamlined Manager Approval / Rejection Portal",
        "Transparent Live Status Tracking for Employees",
        "Centralized MySQL Relational Database Persistence",
        "Robust RESTful Backend APIs built with Spring Boot"
    ]
    for pt in sol_points:
        p = tf_l3.add_paragraph()
        p.space_before = Pt(6)
        p.text = f"✔  {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK

    # Right Card: Clean Workflow Diagram
    c3_right = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.5), Inches(4.9), Inches(4.3))
    c3_right.fill.solid()
    c3_right.fill.fore_color.rgb = CARD_BG
    c3_right.line.color.rgb = CARD_BORDER
    tf_r3 = c3_right.text_frame
    tf_r3.margin_left = Inches(0.3)
    tf_r3.margin_top = Inches(0.3)
    
    p = tf_r3.paragraphs[0]
    p.text = "DIGITAL WORKFLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    wf_steps = ["Employee Login", "Apply Leave Online", "Manager Review", "Approve / Reject", "Updated Status Sync"]
    y_pos = 3.2
    for idx, step in enumerate(wf_steps):
        st_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(y_pos), Inches(4.3), Inches(0.42))
        st_box.fill.solid()
        st_box.fill.fore_color.rgb = PILL_BG
        st_box.line.color.rgb = ACCENT_BLUE
        tf_s = st_box.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.alignment = PP_ALIGN.CENTER
        p_s.text = f"{idx+1}. {step}"
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = NAVY
        y_pos += 0.52
        if idx < len(wf_steps) - 1:
            arr = s3.shapes.add_textbox(Inches(9.8), Inches(y_pos - 0.18), Inches(0.5), Inches(0.2))
            arr.text_frame.paragraphs[0].text = "↓"
            arr.text_frame.paragraphs[0].font.bold = True
            arr.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    # -------------------------------------------------------------
    # SLIDE 4 — KEY FEATURES
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "Key Features", num=4)

    # 3 Column Cards
    col_width = Inches(3.7)
    col_gap = Inches(0.3)
    x_left = Inches(0.8)

    # Employee Column
    c4_1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_left, Inches(1.6), col_width, Inches(5.1))
    c4_1.fill.solid()
    c4_1.fill.fore_color.rgb = CARD_BG
    c4_1.line.color.rgb = CARD_BORDER
    tf4_1 = c4_1.text_frame
    tf4_1.margin_left = Inches(0.3)
    tf4_1.margin_top = Inches(0.3)
    p = tf4_1.paragraphs[0]
    p.text = "👤 EMPLOYEE FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    emp_feats = [
        "Secure User Authentication",
        "Personalized Dashboard",
        "Real-Time Leave Balances",
        "Online Leave Application",
        "Comprehensive Leave History",
        "Track Live Approval Status"
    ]
    for ef in emp_feats:
        p = tf4_1.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {ef}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # Manager Column
    x_left += col_width + col_gap
    c4_2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_left, Inches(1.6), col_width, Inches(5.1))
    c4_2.fill.solid()
    c4_2.fill.fore_color.rgb = CARD_BG
    c4_2.line.color.rgb = CARD_BORDER
    tf4_2 = c4_2.text_frame
    tf4_2.margin_left = Inches(0.3)
    tf4_2.margin_top = Inches(0.3)
    p = tf4_2.paragraphs[0]
    p.text = "👔 MANAGER FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    mgr_feats = [
        "View All Pending Requests",
        "Review Detailed Applications",
        "One-Click Approve / Reject",
        "Add Manager Comments",
        "Monitor Department Leaves",
        "Employee Audit Trail"
    ]
    for mf in mgr_feats:
        p = tf4_2.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {mf}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # System Column
    x_left += col_width + col_gap
    c4_3 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_left, Inches(1.6), col_width, Inches(5.1))
    c4_3.fill.solid()
    c4_3.fill.fore_color.rgb = CARD_BG
    c4_3.line.color.rgb = CARD_BORDER
    tf4_3 = c4_3.text_frame
    tf4_3.margin_left = Inches(0.3)
    tf4_3.margin_top = Inches(0.3)
    p = tf4_3.paragraphs[0]
    p.text = "⚙️ SYSTEM FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sys_feats = [
        "REST API Architecture",
        "MySQL Relational Database",
        "Robust Error Handling",
        "Strict Input API Validation",
        "Playwright E2E Testing",
        "CORS & Security Config"
    ]
    for sf in sys_feats:
        p = tf4_3.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {sf}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 5 — SYSTEM ARCHITECTURE
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "System Architecture", num=5)

    # Architecture Diagram Container
    arch_bg = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.3))
    arch_bg.fill.solid()
    arch_bg.fill.fore_color.rgb = CARD_BG
    arch_bg.line.color.rgb = CARD_BORDER

    # Title inside diagram
    t_arch = arch_bg.text_frame
    t_arch.margin_top = Inches(0.2)
    p_arch = t_arch.paragraphs[0]
    p_arch.alignment = PP_ALIGN.CENTER
    p_arch.text = "CAMPUSLOOP HIGH-LEVEL ARCHITECTURE"
    p_arch.font.size = Pt(14)
    p_arch.font.bold = True
    p_arch.font.color.rgb = NAVY

    # Box 1: React Frontend
    b1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.3), Inches(3.2), Inches(2.8))
    b1.fill.solid()
    b1.fill.fore_color.rgb = PILL_BG
    b1.line.color.rgb = ACCENT_BLUE
    tf1 = b1.text_frame
    tf1.margin_left = Inches(0.2)
    p = tf1.paragraphs[0]
    p.text = "REACT FRONTEND\n(Client Layer)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    f_items = ["• React.js (JSX)", "• Vite Build Tool", "• React Router", "• Auth Context", "• Dashboard & Forms"]
    for fi in f_items:
        p = tf1.add_paragraph()
        p.space_before = Pt(4)
        p.text = fi
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK

    # Arrow 1
    a1 = s5.shapes.add_textbox(Inches(4.45), Inches(3.4), Inches(0.9), Inches(0.4))
    a1.text_frame.paragraphs[0].text = "⇄ REST\nAPIs"
    a1.text_frame.paragraphs[0].font.size = Pt(10)
    a1.text_frame.paragraphs[0].font.bold = True
    a1.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    # Box 2: Spring Boot Backend
    b2 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.3), Inches(3.4), Inches(2.8))
    b2.fill.solid()
    b2.fill.fore_color.rgb = PILL_BG
    b2.line.color.rgb = DARK_BLUE
    tf2 = b2.text_frame
    tf2.margin_left = Inches(0.2)
    p = tf2.paragraphs[0]
    p.text = "SPRING BOOT BACKEND\n(Business Logic)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    b_items = ["• REST Controllers", "• Service Layer", "• Repository Layer", "• Spring Security", "• DTOs & Validation"]
    for bi in b_items:
        p = tf2.add_paragraph()
        p.space_before = Pt(4)
        p.text = bi
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK

    # Arrow 2
    a2 = s5.shapes.add_textbox(Inches(8.75), Inches(3.4), Inches(0.9), Inches(0.4))
    a2.text_frame.paragraphs[0].text = "⇄ JPA /\nJDBC"
    a2.text_frame.paragraphs[0].font.size = Pt(10)
    a2.text_frame.paragraphs[0].font.bold = True
    a2.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    # Box 3: MySQL Database
    b3 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.6), Inches(2.3), Inches(2.5), Inches(2.8))
    b3.fill.solid()
    b3.fill.fore_color.rgb = PILL_BG
    b3.line.color.rgb = NAVY
    tf3 = b3.text_frame
    tf3.margin_left = Inches(0.2)
    p = tf3.paragraphs[0]
    p.text = "MYSQL DB\n(Persistence)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY

    d_items = ["• User Table", "• Employee Table", "• LeaveRequest Table", "• Relational Keys"]
    for di in d_items:
        p = tf3.add_paragraph()
        p.space_before = Pt(4)
        p.text = di
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK

    # Bottom Explanation Callout
    exp_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = PILL_BG
    exp_box.line.color.rgb = ACCENT_BLUE
    tf_exp = exp_box.text_frame
    p_exp = tf_exp.paragraphs[0]
    p_exp.alignment = PP_ALIGN.CENTER
    p_exp.text = "Architecture Summary: The React frontend communicates with Spring Boot REST APIs, which process business logic and persist employee & leave information in MySQL. API quality is verified using Postman and Playwright."
    p_exp.font.size = Pt(11)
    p_exp.font.bold = True
    p_exp.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 6 — EMPLOYEE → MANAGER LEAVE WORKFLOW
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Employee → Manager Leave Workflow", num=6)

    # 6 Steps Grid (2 rows x 3 columns)
    steps_data = [
        ("STEP 1", "Employee Login", "Employee authenticates using secure email and password."),
        ("STEP 2", "Apply Leave", "Selects Leave Type, Start/End Dates, Number of Days & Reason."),
        ("STEP 3", "Request Stored", "Leave request is persisted in MySQL database with PENDING status."),
        ("STEP 4", "Manager Reviews", "Manager views pending leave requests in their portal."),
        ("STEP 5", "Manager Decision", "Manager evaluates the request and selects APPROVE or REJECT."),
        ("STEP 6", "Status Updated", "Employee views updated status and updated balance on dashboard.")
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.3)
    
    positions = [
        (Inches(0.8), Inches(1.6)), (Inches(4.8), Inches(1.6)), (Inches(8.8), Inches(1.6)),
        (Inches(0.8), Inches(4.3)), (Inches(4.8), Inches(4.3)), (Inches(8.8), Inches(4.3))
    ]

    for idx, (step_num, title, desc) in enumerate(steps_data):
        x, y = positions[idx]
        scard = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        scard.fill.solid()
        scard.fill.fore_color.rgb = CARD_BG
        scard.line.color.rgb = CARD_BORDER
        
        tfs = scard.text_frame
        tfs.margin_left = Inches(0.2)
        tfs.margin_top = Inches(0.2)
        
        # Step Badge
        p_num = tfs.paragraphs[0]
        p_num.text = step_num
        p_num.font.size = Pt(10)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_BLUE
        
        p_title = tfs.add_paragraph()
        p_title.space_before = Pt(4)
        p_title.text = title
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        
        p_desc = tfs.add_paragraph()
        p_desc.space_before = Pt(6)
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 7 — AI-ASSISTED DEVELOPMENT
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "AI-Assisted Development", num=7)

    # Left Card: Assistance points
    c7_left = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.1))
    c7_left.fill.solid()
    c7_left.fill.fore_color.rgb = CARD_BG
    c7_left.line.color.rgb = CARD_BORDER
    tf7_l = c7_left.text_frame
    tf7_l.margin_left = Inches(0.3)
    tf7_l.margin_top = Inches(0.3)

    p = tf7_l.paragraphs[0]
    p.text = "AI ACCELERATION & DEVELOPER VALIDATION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    ai_points = [
        "Assisted in generating and refining boilerplate code",
        "Accelerated Spring Boot API controller & service layer debugging",
        "Assisted with React UI component layout & state management",
        "Helped identify and resolve CORS & API integration issues",
        "Assisted in creating comprehensive test scenarios",
        "Helped optimize backend error handling and validation logic",
        "Supported Playwright end-to-end test automation setup",
        "Developer reviewed, modified, and validated all generated code"
    ]
    for ap in ai_points:
        p = tf7_l.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"•  {ap}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # Right Card: Process Diagram & Transparency Statement
    c7_right = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.6), Inches(4.2), Inches(5.1))
    c7_right.fill.solid()
    c7_right.fill.fore_color.rgb = CARD_BG
    c7_right.line.color.rgb = CARD_BORDER
    tf7_r = c7_right.text_frame
    tf7_r.margin_left = Inches(0.3)
    tf7_r.margin_right = Inches(0.3)
    tf7_r.margin_top = Inches(0.3)

    p = tf7_r.paragraphs[0]
    p.text = "DEVELOPMENT WORKFLOW"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # Workflow steps
    w7_steps = ["Developer Intent", "AI Assistance Prompt", "Code Generation", "Manual Debug & Refine", "Validation & Testing"]
    for idx, ws in enumerate(w7_steps):
        p = tf7_r.add_paragraph()
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        p.text = f"[{idx+1}]  {ws}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

    p_disc = tf7_r.add_paragraph()
    p_disc.space_before = Pt(20)
    p_disc.text = "Note: AI tools were used purely as development assistance. All application logic, schema, and API integrations were thoroughly verified and tested by the developer."
    p_disc.font.size = Pt(10)
    p_disc.font.italic = True
    p_disc.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 8 — API TESTING & TEST AUTOMATION
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "API Testing & Test Automation", num=8)

    # Left Column: Tested Endpoints
    c8_left = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1))
    c8_left.fill.solid()
    c8_left.fill.fore_color.rgb = CARD_BG
    c8_left.line.color.rgb = CARD_BORDER
    tf8_l = c8_left.text_frame
    tf8_l.margin_left = Inches(0.3)
    tf8_l.margin_top = Inches(0.3)

    p = tf8_l.paragraphs[0]
    p.text = "VALIDATED REST API ENDPOINTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    apis = [
        ("POST /api/auth/login", "Employee authentication & JWT session"),
        ("POST /api/leaves", "Create a new leave request"),
        ("GET /api/leaves", "Retrieve all leave applications"),
        ("GET /api/leaves/employee/{id}", "Retrieve employee specific history"),
        ("PUT /api/leaves/{id}/status", "Update leave status (Approve/Reject)")
    ]
    for ep, desc in apis:
        p = tf8_l.add_paragraph()
        p.space_before = Pt(10)
        p.text = ep
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        p2 = tf8_l.add_paragraph()
        p2.text = f"Purpose: {desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_DARK

    # Right Column: Postman Screenshots
    # Screenshot 1: 200 OK Login
    if os.path.exists("scratch_assets/postman_200_login.png"):
        s8.shapes.add_picture("scratch_assets/postman_200_login.png", Inches(6.7), Inches(1.6), Inches(5.8), Inches(2.4))

    # Screenshot 2: 201 Created Leave
    if os.path.exists("scratch_assets/postman_201_create.png"):
        s8.shapes.add_picture("scratch_assets/postman_201_create.png", Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.4))

    # -------------------------------------------------------------
    # SLIDE 9 — DELIBERATE FAILURE → FIX → PASS
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Test Validation: Failure → Fix → Pass", num=9)

    # 3 Stage Top Container
    stage_w = Inches(3.7)
    gap = Inches(0.3)
    x = Inches(0.8)

    # Stage 1: Deliberate Failure
    s1_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), stage_w, Inches(2.2))
    s1_box.fill.solid()
    s1_box.fill.fore_color.rgb = RGBColor(254, 242, 242)
    s1_box.line.color.rgb = RGBColor(252, 165, 165)
    t1 = s1_box.text_frame
    t1.margin_left = Inches(0.2)
    p = t1.paragraphs[0]
    p.text = "🔴 1. DELIBERATE FAILURE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DANGER_RED

    p = t1.add_paragraph()
    p.space_before = Pt(6)
    p.text = "Login tested with wrong password."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    p = t1.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Result: 401 Unauthorized\n'Invalid email or password'"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DANGER_RED

    # Stage 2: Fix
    x += stage_w + gap
    s2_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), stage_w, Inches(2.2))
    s2_box.fill.solid()
    s2_box.fill.fore_color.rgb = PILL_BG
    s2_box.line.color.rgb = ACCENT_BLUE
    t2 = s2_box.text_frame
    t2.margin_left = Inches(0.2)
    p = t2.paragraphs[0]
    p.text = "🛠 2. FIX"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p = t2.add_paragraph()
    p.space_before = Pt(6)
    p.text = "Correct employee credentials provided."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    p = t2.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Action: Request payload updated & re-sent to /api/auth/login"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Stage 3: Pass
    x += stage_w + gap
    s3_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), stage_w, Inches(2.2))
    s3_box.fill.solid()
    s3_box.fill.fore_color.rgb = RGBColor(236, 253, 245)
    s3_box.line.color.rgb = RGBColor(110, 231, 183)
    t3 = s3_box.text_frame
    t3.margin_left = Inches(0.2)
    p = t3.paragraphs[0]
    p.text = "🟢 3. PASS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN

    p = t3.add_paragraph()
    p.space_before = Pt(6)
    p.text = "Valid authentication request."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    p = t3.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Result: 200 OK\n'Authentication successful'"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN

    # Postman Screenshots Comparison
    # Left: 401 Unauthorized
    if os.path.exists("scratch_assets/postman_401_fail.png"):
        s9.shapes.add_picture("scratch_assets/postman_401_fail.png", Inches(0.8), Inches(4.1), Inches(5.7), Inches(2.6))

    # Right: 200 OK
    if os.path.exists("scratch_assets/postman_200_login.png"):
        s9.shapes.add_picture("scratch_assets/postman_200_login.png", Inches(6.8), Inches(4.1), Inches(5.7), Inches(2.6))

    # -------------------------------------------------------------
    # SLIDE 10 — CONCLUSION + DEMO
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_header(s10, "Conclusion & Live Demo", num=10)

    # Left Column: Key Outcomes
    c10_left = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.2))
    c10_left.fill.solid()
    c10_left.fill.fore_color.rgb = CARD_BG
    c10_left.line.color.rgb = CARD_BORDER
    tf10_l = c10_left.text_frame
    tf10_l.margin_left = Inches(0.3)
    tf10_l.margin_top = Inches(0.3)

    p = tf10_l.paragraphs[0]
    p.text = "PROJECT OUTCOMES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    outcomes = [
        "Simplifies and digitizes employee leave application",
        "Reduces manual administrative paperwork and delays",
        "Provides transparent real-time status & leave balances",
        "Improves manager approval and decision workflow",
        "Robust Spring Boot REST architecture & MySQL persistence",
        "Validated with Playwright automated tests & Postman"
    ]
    for oc in outcomes:
        p = tf10_l.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"✔  {oc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # Right Column: Demo Flow
    c10_right = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.2))
    c10_right.fill.solid()
    c10_right.fill.fore_color.rgb = CARD_BG
    c10_right.line.color.rgb = CARD_BORDER
    tf10_r = c10_right.text_frame
    tf10_r.margin_left = Inches(0.3)
    tf10_r.margin_top = Inches(0.3)

    p = tf10_r.paragraphs[0]
    p.text = "LIVE DEMO FLOW"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    demo_steps = [
        "1. Employee Login & Dashboard",
        "2. View Leave Balances",
        "3. Submit New Leave Request",
        "4. View Pending Application",
        "5. Manager Reviews Request",
        "6. Manager Approve / Reject Action",
        "7. Employee View Updated Status",
        "8. Demonstrate Postman API Validation"
    ]
    for ds in demo_steps:
        p = tf10_r.add_paragraph()
        p.space_before = Pt(4)
        p.text = ds
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK

    # Bottom Banner: Tagline & Thank You
    b10 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9))
    b10.fill.solid()
    b10.fill.fore_color.rgb = NAVY
    b10.line.fill.background()
    tf_b10 = b10.text_frame
    p_tag = tf_b10.paragraphs[0]
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.text = "CampusLoop — Simplifying Employee Leave Management Through Technology"
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(255, 255, 255)

    p_ty = tf_b10.add_paragraph()
    p_ty.space_before = Pt(4)
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.text = "THANK YOU!  |  Q & A SESSION"
    p_ty.font.size = Pt(11)
    p_ty.font.bold = True
    p_ty.font.color.rgb = RGBColor(147, 197, 253)

    output_path = "CampusLoop_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    build_deck()
