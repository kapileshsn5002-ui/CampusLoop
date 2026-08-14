import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_perfect_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # High-Contrast Colors
    BLACK_TEXT = RGBColor(15, 23, 42)      # #0F172A Very Dark Slate / Black
    NAVY_TITLE = RGBColor(30, 58, 138)    # #1E3A8A Deep Blue
    ACCENT_BLUE = RGBColor(37, 99, 235)   # #2563EB Bright Blue
    MUTED_TEXT = RGBColor(71, 85, 105)    # #475569 Dark Muted Slate
    
    WHITE_BG = RGBColor(255, 255, 255)    # #FFFFFF Pure White
    SLATE_BG = RGBColor(248, 250, 252)    # #F8FAFC Light Slate
    CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF Card BG
    CARD_BORDER = RGBColor(148, 163, 184) # #94A3B8 Strong Border
    
    BLUE_BOX_BG = RGBColor(219, 234, 254)  # #DBEAFE Light Blue Fill
    BLUE_BOX_BORDER = RGBColor(37, 99, 235)# #2563EB
    
    RED_BOX_BG = RGBColor(254, 226, 226)   # #FEE2E2 Light Red
    RED_TEXT = RGBColor(185, 28, 28)       # #B91C1C Dark Red
    RED_BORDER = RGBColor(239, 68, 68)    # #EF4444
    
    GREEN_BOX_BG = RGBColor(209, 250, 229) # #D1FAE5 Light Green
    GREEN_TEXT = RGBColor(4, 120, 87)      # #047857 Dark Green
    GREEN_BORDER = RGBColor(16, 185, 129)  # #10B981

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = SLATE_BG
        bg.line.fill.background()

    def add_header(slide, title_text, slide_num):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.3))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "CAMPUSLOOP PRESENTATION"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.65), Inches(10.0), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_TITLE

        # Line Divider
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

        # Footer Slide Number
        foot = slide.shapes.add_textbox(Inches(11.0), Inches(7.0), Inches(1.5), Inches(0.4))
        p_f = foot.text_frame.paragraphs[0]
        p_f.alignment = PP_ALIGN.RIGHT
        p_f.text = f"Slide {slide_num} / 10"
        p_f.font.size = Pt(12)
        p_f.font.bold = True
        p_f.font.color.rgb = BLACK_TEXT

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # -------------------------------------------------------------
    # SLIDE 1 — TITLE
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_bg(s1)

    # Accent Stripe
    stripe = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_BLUE
    stripe.line.fill.background()

    # Academic Badge Box
    badge = add_card(s1, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.45), BLUE_BOX_BG, BLUE_BOX_BORDER)
    p = badge.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "ACADEMIC PROJECT PRESENTATION (2025–2026)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # Main Title
    t1 = s1.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(7.5), Inches(1.0))
    p = t1.text_frame.paragraphs[0]
    p.text = "CampusLoop"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # Subtitle
    st1 = s1.shapes.add_textbox(Inches(0.75), Inches(2.4), Inches(7.5), Inches(0.5))
    p = st1.text_frame.paragraphs[0]
    p.text = "Employee Leave Management System"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLACK_TEXT

    # Tagline
    tag1 = s1.shapes.add_textbox(Inches(0.75), Inches(3.0), Inches(7.5), Inches(0.5))
    p = tag1.text_frame.paragraphs[0]
    p.text = "Web-Based Employee Leave Management & Approval Platform"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MUTED_TEXT

    # Tech Stack Box
    tech1 = add_card(s1, Inches(0.8), Inches(3.7), Inches(7.2), Inches(0.6), BLUE_BOX_BG, BLUE_BOX_BORDER)
    p = tech1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "React.js  |  Spring Boot  |  MySQL  |  REST API  |  Playwright"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # Student Info Right Card
    card_student = add_card(s1, Inches(8.4), Inches(0.8), Inches(4.1), Inches(5.8), WHITE_BG, CARD_BORDER)
    
    # We place a dedicated text box over the student card to guarantee black text
    tb_student = s1.shapes.add_textbox(Inches(8.6), Inches(1.0), Inches(3.7), Inches(5.4))
    tf_s = tb_student.text_frame
    tf_s.word_wrap = True
    
    p = tf_s.paragraphs[0]
    p.text = "STUDENT DETAILS"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    student_data = [
        ("Student Name:", "[Student Name]"),
        ("Register Number:", "[Register Number]"),
        ("Department:", "[Department]"),
        ("College Name:", "[College Name]"),
        ("Academic Year:", "2025–2026")
    ]
    for lbl, val in student_data:
        p1 = tf_s.add_paragraph()
        p1.space_before = Pt(14)
        p1.text = lbl
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = MUTED_TEXT

        p2 = tf_s.add_paragraph()
        p2.text = val
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = BLACK_TEXT

    # Footer Slide 1
    f1 = s1.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(3.0), Inches(0.4))
    p = f1.text_frame.paragraphs[0]
    p.text = "Slide 1 / 10"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK_TEXT

    # -------------------------------------------------------------
    # SLIDE 2 — PROBLEM STATEMENT
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_bg(s2)
    add_header(s2, "Problem Statement", 2)

    # Left Card
    add_card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.1), WHITE_BG, CARD_BORDER)
    tb2_l = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(3.7))
    tf2_l = tb2_l.text_frame
    tf2_l.word_wrap = True

    p = tf2_l.paragraphs[0]
    p.text = "TRADITIONAL LEAVE MANAGEMENT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED_TEXT

    prob_bullets = [
        "Manual paper-based leave applications",
        "Informal & unorganized approval workflows",
        "Difficulty tracking real-time leave balances",
        "Lack of centralized employee record-keeping",
        "Delayed manager responses & approval bottlenecks",
        "Limited visibility of application status for employees",
        "Increased risk of data-entry & accounting errors"
    ]
    for b in prob_bullets:
        p = tf2_l.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"•  {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Right Card
    add_card(s2, Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.1), WHITE_BG, CARD_BORDER)
    tb2_r = s2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(3.7))
    tf2_r = tb2_r.text_frame
    tf2_r.word_wrap = True

    p = tf2_r.paragraphs[0]
    p.text = "CORE PROBLEM DEFINITION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    p = tf2_r.add_paragraph()
    p.space_before = Pt(14)
    p.text = "Employees and managers require a centralized digital web platform where leave requests can be submitted, tracked, reviewed, and approved efficiently without administrative friction."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK_TEXT

    # Traditional Flow Box
    add_card(s2, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.9), RED_BOX_BG, RED_BORDER)
    tb2_fl = s2.shapes.add_textbox(Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6))
    p = tb2_fl.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Traditional Flow (Flawed):  Employee  ➔  Manual Paper Form  ➔  Manager Desk  ➔  Delayed Approval  ➔  Manual Records"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RED_TEXT

    # -------------------------------------------------------------
    # SLIDE 3 — PROPOSED SOLUTION
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_bg(s3)
    add_header(s3, "Proposed Solution", 3)

    # Top Banner
    add_card(s3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.85), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb3_top = s3.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.65))
    p = tb3_top.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CampusLoop is a web-based employee leave management system designed to digitize, automate, and simplify the complete organizational leave workflow."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # Left Card
    add_card(s3, Inches(0.8), Inches(2.55), Inches(6.5), Inches(4.35), WHITE_BG, CARD_BORDER)
    tb3_l = s3.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(6.1), Inches(4.0))
    tf3_l = tb3_l.text_frame
    tf3_l.word_wrap = True

    p = tf3_l.paragraphs[0]
    p.text = "KEY SOLUTION ADVANTAGES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    sol_bullets = [
        "Secure Employee & Manager Authentication",
        "Interactive Employee Dashboard with Leave History",
        "Online Application Submission with Instant Validation",
        "Real-Time Leave Balance Calculation & Tracking",
        "Streamlined Manager Approval / Rejection Portal",
        "Transparent Live Status Tracking for Employees",
        "Centralized MySQL Relational Database Persistence",
        "Robust RESTful Backend APIs built with Spring Boot"
    ]
    for b in sol_bullets:
        p = tf3_l.add_paragraph()
        p.space_before = Pt(6)
        p.text = f"✔  {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Right Card: Workflow
    add_card(s3, Inches(7.6), Inches(2.55), Inches(4.9), Inches(4.35), WHITE_BG, CARD_BORDER)
    tb3_r_hdr = s3.shapes.add_textbox(Inches(7.8), Inches(2.7), Inches(4.5), Inches(0.4))
    p = tb3_r_hdr.text_frame.paragraphs[0]
    p.text = "DIGITAL WORKFLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    wf_steps = ["1. Employee Login", "2. Apply Leave Online", "3. Manager Review", "4. Approve / Reject", "5. Updated Status Sync"]
    y_pos = 3.2
    for step in wf_steps:
        box = add_card(s3, Inches(7.9), Inches(y_pos), Inches(4.3), Inches(0.45), BLUE_BOX_BG, BLUE_BOX_BORDER)
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = step
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY_TITLE
        y_pos += 0.55

    # -------------------------------------------------------------
    # SLIDE 4 — KEY FEATURES
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_bg(s4)
    add_header(s4, "Key Features Grid", 4)

    col_w = Inches(3.7)
    col_gap = Inches(0.3)
    x = Inches(0.8)

    # Column 1: Employee
    add_card(s4, x, Inches(1.6), col_w, Inches(5.1), WHITE_BG, CARD_BORDER)
    tb4_1 = s4.shapes.add_textbox(x + Inches(0.2), Inches(1.8), col_w - Inches(0.4), Inches(4.7))
    tf4_1 = tb4_1.text_frame
    tf4_1.word_wrap = True
    p = tf4_1.paragraphs[0]
    p.text = "👤 EMPLOYEE FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    ef_list = ["Secure User Authentication", "Personalized Dashboard", "Real-Time Leave Balances", "Online Leave Application", "Comprehensive Leave History", "Track Live Approval Status"]
    for item in ef_list:
        p = tf4_1.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Column 2: Manager
    x += col_w + col_gap
    add_card(s4, x, Inches(1.6), col_w, Inches(5.1), WHITE_BG, CARD_BORDER)
    tb4_2 = s4.shapes.add_textbox(x + Inches(0.2), Inches(1.8), col_w - Inches(0.4), Inches(4.7))
    tf4_2 = tb4_2.text_frame
    tf4_2.word_wrap = True
    p = tf4_2.paragraphs[0]
    p.text = "👔 MANAGER FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    mf_list = ["View Pending Requests", "Review Applications", "One-Click Approve/Reject", "Add Manager Comments", "Department Balance Audit", "Action History Trail"]
    for item in mf_list:
        p = tf4_2.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Column 3: System
    x += col_w + col_gap
    add_card(s4, x, Inches(1.6), col_w, Inches(5.1), WHITE_BG, CARD_BORDER)
    tb4_3 = s4.shapes.add_textbox(x + Inches(0.2), Inches(1.8), col_w - Inches(0.4), Inches(4.7))
    tf4_3 = tb4_3.text_frame
    tf4_3.word_wrap = True
    p = tf4_3.paragraphs[0]
    p.text = "⚙️ SYSTEM FEATURES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    sf_list = ["REST API Integration", "MySQL Relational DB", "Robust Error Handling", "Strict Payload Validation", "Automated Playwright Tests", "CORS Security Setup"]
    for item in sf_list:
        p = tf4_3.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # -------------------------------------------------------------
    # SLIDE 5 — SYSTEM ARCHITECTURE
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_bg(s5)
    add_header(s5, "System Architecture", 5)

    add_card(s5, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.3), WHITE_BG, CARD_BORDER)
    tb5_hdr = s5.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.3), Inches(0.4))
    p = tb5_hdr.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CAMPUSLOOP HIGH-LEVEL ARCHITECTURE"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # Box 1: React Frontend
    add_card(s5, Inches(1.2), Inches(2.3), Inches(3.2), Inches(2.8), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb5_1 = s5.shapes.add_textbox(Inches(1.3), Inches(2.4), Inches(3.0), Inches(2.6))
    tf5_1 = tb5_1.text_frame
    tf5_1.word_wrap = True
    p = tf5_1.paragraphs[0]
    p.text = "REACT FRONTEND\n(Client Layer)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    f_items = ["• React.js (JSX)", "• Vite Build Tool", "• React Router", "• Auth Context", "• Dashboard & Forms"]
    for fi in f_items:
        p = tf5_1.add_paragraph()
        p.space_before = Pt(4)
        p.text = fi
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK_TEXT

    # Arrow 1
    tb_a1 = s5.shapes.add_textbox(Inches(4.45), Inches(3.4), Inches(0.9), Inches(0.5))
    p = tb_a1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "⇄ REST\nAPIs"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # Box 2: Spring Boot Backend
    add_card(s5, Inches(5.3), Inches(2.3), Inches(3.4), Inches(2.8), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb5_2 = s5.shapes.add_textbox(Inches(5.4), Inches(2.4), Inches(3.2), Inches(2.6))
    tf5_2 = tb5_2.text_frame
    tf5_2.word_wrap = True
    p = tf5_2.paragraphs[0]
    p.text = "SPRING BOOT BACKEND\n(Business Logic)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    b_items = ["• REST Controllers", "• Service Layer", "• Repository Layer", "• Spring Security", "• DTOs & Validation"]
    for bi in b_items:
        p = tf5_2.add_paragraph()
        p.space_before = Pt(4)
        p.text = bi
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK_TEXT

    # Arrow 2
    tb_a2 = s5.shapes.add_textbox(Inches(8.75), Inches(3.4), Inches(0.9), Inches(0.5))
    p = tb_a2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "⇄ JPA /\nJDBC"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # Box 3: MySQL DB
    add_card(s5, Inches(9.6), Inches(2.3), Inches(2.5), Inches(2.8), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb5_3 = s5.shapes.add_textbox(Inches(9.7), Inches(2.4), Inches(2.3), Inches(2.6))
    tf5_3 = tb5_3.text_frame
    tf5_3.word_wrap = True
    p = tf5_3.paragraphs[0]
    p.text = "MYSQL DB\n(Persistence)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    d_items = ["• User Table", "• Employee Table", "• LeaveRequest Table", "• Relational Keys"]
    for di in d_items:
        p = tf5_3.add_paragraph()
        p.space_before = Pt(4)
        p.text = di
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK_TEXT

    # Bottom Explanation
    add_card(s5, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.85), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb5_exp = s5.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.65))
    p = tb5_exp.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Architecture Summary: The React frontend communicates with Spring Boot REST APIs, which process business logic and persist employee & leave information in MySQL. API quality is verified using Postman and Playwright."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLACK_TEXT

    # -------------------------------------------------------------
    # SLIDE 6 — EMPLOYEE → MANAGER LEAVE WORKFLOW
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_bg(s6)
    add_header(s6, "Employee → Manager Leave Workflow", 6)

    steps_data = [
        ("STEP 1", "Employee Login", "Employee authenticates using secure email and password."),
        ("STEP 2", "Apply Leave", "Selects Leave Type, Start/End Dates, Number of Days & Reason."),
        ("STEP 3", "Request Stored", "Leave request is persisted in MySQL database with PENDING status."),
        ("STEP 4", "Manager Reviews", "Manager views pending leave requests in their portal."),
        ("STEP 5", "Manager Decision", "Manager evaluates the request and selects APPROVE or REJECT."),
        ("STEP 6", "Status Updated", "Employee views updated status and updated balance on dashboard.")
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.3)
    positions = [
        (Inches(0.8), Inches(1.6)), (Inches(4.8), Inches(1.6)), (Inches(8.8), Inches(1.6)),
        (Inches(0.8), Inches(4.3)), (Inches(4.8), Inches(4.3)), (Inches(8.8), Inches(4.3))
    ]

    for idx, (step_num, title, desc) in enumerate(steps_data):
        x_p, y_p = positions[idx]
        add_card(s6, x_p, y_p, card_w, card_h, WHITE_BG, CARD_BORDER)
        
        tb = s6.shapes.add_textbox(x_p + Inches(0.2), y_p + Inches(0.2), card_w - Inches(0.4), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = step_num
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY_TITLE

        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = BLACK_TEXT

    # -------------------------------------------------------------
    # SLIDE 7 — AI-ASSISTED DEVELOPMENT
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_bg(s7)
    add_header(s7, "AI-Assisted Development", 7)

    # Left Card
    add_card(s7, Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.1), WHITE_BG, CARD_BORDER)
    tb7_l = s7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(6.8), Inches(4.7))
    tf7_l = tb7_l.text_frame
    tf7_l.word_wrap = True

    p = tf7_l.paragraphs[0]
    p.text = "AI ACCELERATION & DEVELOPER VALIDATION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    ai_pts = [
        "Assisted in generating and refining boilerplate code",
        "Accelerated Spring Boot API controller & service layer debugging",
        "Assisted with React UI component layout & state management",
        "Helped identify and resolve CORS & API integration issues",
        "Assisted in creating comprehensive test scenarios",
        "Helped optimize backend error handling and validation logic",
        "Supported Playwright end-to-end test automation setup",
        "Developer reviewed, modified, and validated all generated code"
    ]
    for item in ai_pts:
        p = tf7_l.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"•  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Right Card
    add_card(s7, Inches(8.3), Inches(1.6), Inches(4.2), Inches(5.1), WHITE_BG, CARD_BORDER)
    tb7_r = s7.shapes.add_textbox(Inches(8.5), Inches(1.8), Inches(3.8), Inches(4.7))
    tf7_r = tb7_r.text_frame
    tf7_r.word_wrap = True

    p = tf7_r.paragraphs[0]
    p.text = "DEVELOPMENT WORKFLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    flow_pts = ["Developer Intent", "AI Assistance Prompt", "Code Generation", "Manual Debug & Refine", "Validation & Testing"]
    for idx, item in enumerate(flow_pts):
        p = tf7_r.add_paragraph()
        p.space_before = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        p.text = f"[{idx+1}]  {item}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY_TITLE

    p = tf7_r.add_paragraph()
    p.space_before = Pt(24)
    p.text = "Note: AI tools were used purely as development assistance. All application logic, schema, and API integrations were thoroughly verified and tested by the developer."
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = MUTED_TEXT

    # -------------------------------------------------------------
    # SLIDE 8 — API TESTING & TEST AUTOMATION
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_bg(s8)
    add_header(s8, "API Testing & Test Automation", 8)

    # Left Column
    add_card(s8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), WHITE_BG, CARD_BORDER)
    tb8_l = s8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.7))
    tf8_l = tb8_l.text_frame
    tf8_l.word_wrap = True

    p = tf8_l.paragraphs[0]
    p.text = "VALIDATED REST API ENDPOINTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    apis = [
        ("POST /api/auth/login", "Employee authentication & JWT session"),
        ("POST /api/leaves", "Create a new leave request"),
        ("GET /api/leaves", "Retrieve all leave applications"),
        ("GET /api/leaves/employee/{id}", "Retrieve employee specific history"),
        ("PUT /api/leaves/{id}/status", "Update leave status (Approve/Reject)")
    ]
    for ep, desc in apis:
        p = tf8_l.add_paragraph()
        p.space_before = Pt(10)
        p.text = ep
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        p2 = tf8_l.add_paragraph()
        p2.text = f"Purpose: {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = BLACK_TEXT

    # Right Screenshots
    if os.path.exists("scratch_assets/postman_200_login.png"):
        s8.shapes.add_picture("scratch_assets/postman_200_login.png", Inches(6.7), Inches(1.6), Inches(5.8), Inches(2.4))

    if os.path.exists("scratch_assets/postman_201_create.png"):
        s8.shapes.add_picture("scratch_assets/postman_201_create.png", Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.4))

    # -------------------------------------------------------------
    # SLIDE 9 — DELIBERATE FAILURE → FIX → PASS
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_bg(s9)
    add_header(s9, "Test Validation: Failure → Fix → Pass", 9)

    sw = Inches(3.7)
    sgap = Inches(0.3)
    sx = Inches(0.8)

    # 1. Deliberate Failure Box
    add_card(s9, sx, Inches(1.6), sw, Inches(2.2), RED_BOX_BG, RED_BORDER)
    tb9_1 = s9.shapes.add_textbox(sx + Inches(0.15), Inches(1.7), sw - Inches(0.3), Inches(2.0))
    tf9_1 = tb9_1.text_frame
    tf9_1.word_wrap = True
    p = tf9_1.paragraphs[0]
    p.text = "🔴 1. DELIBERATE FAILURE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RED_TEXT

    p = tf9_1.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Login tested with wrong password."
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK_TEXT

    p = tf9_1.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Result: 401 Unauthorized\n'Invalid email or password'"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED_TEXT

    # 2. Fix Box
    sx += sw + sgap
    add_card(s9, sx, Inches(1.6), sw, Inches(2.2), BLUE_BOX_BG, BLUE_BOX_BORDER)
    tb9_2 = s9.shapes.add_textbox(sx + Inches(0.15), Inches(1.7), sw - Inches(0.3), Inches(2.0))
    tf9_2 = tb9_2.text_frame
    tf9_2.word_wrap = True
    p = tf9_2.paragraphs[0]
    p.text = "🛠 2. FIX"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    p = tf9_2.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Correct employee credentials provided."
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK_TEXT

    p = tf9_2.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Action: Request payload updated & re-sent to /api/auth/login"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    # 3. Pass Box
    sx += sw + sgap
    add_card(s9, sx, Inches(1.6), sw, Inches(2.2), GREEN_BOX_BG, GREEN_BORDER)
    tb9_3 = s9.shapes.add_textbox(sx + Inches(0.15), Inches(1.7), sw - Inches(0.3), Inches(2.0))
    tf9_3 = tb9_3.text_frame
    tf9_3.word_wrap = True
    p = tf9_3.paragraphs[0]
    p.text = "🟢 3. PASS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GREEN_TEXT

    p = tf9_3.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Valid authentication request."
    p.font.size = Pt(11)
    p.font.color.rgb = BLACK_TEXT

    p = tf9_3.add_paragraph()
    p.space_before = Pt(4)
    p.text = "Result: 200 OK\n'Authentication successful'"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN_TEXT

    # Screenshots Comparison
    if os.path.exists("scratch_assets/postman_401_fail.png"):
        s9.shapes.add_picture("scratch_assets/postman_401_fail.png", Inches(0.8), Inches(4.1), Inches(5.7), Inches(2.6))

    if os.path.exists("scratch_assets/postman_200_login.png"):
        s9.shapes.add_picture("scratch_assets/postman_200_login.png", Inches(6.8), Inches(4.1), Inches(5.7), Inches(2.6))

    # -------------------------------------------------------------
    # SLIDE 10 — CONCLUSION + DEMO
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_bg(s10)
    add_header(s10, "Conclusion & Live Demo", 10)

    # Left Card
    add_card(s10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.2), WHITE_BG, CARD_BORDER)
    tb10_l = s10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(3.8))
    tf10_l = tb10_l.text_frame
    tf10_l.word_wrap = True

    p = tf10_l.paragraphs[0]
    p.text = "PROJECT OUTCOMES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE

    outcomes = [
        "Simplifies and digitizes employee leave application",
        "Reduces manual administrative paperwork and delays",
        "Provides transparent real-time status & leave balances",
        "Improves manager approval and decision workflow",
        "Robust Spring Boot REST architecture & MySQL persistence",
        "Validated with Playwright automated tests & Postman"
    ]
    for oc in outcomes:
        p = tf10_l.add_paragraph()
        p.space_before = Pt(8)
        p.text = f"✔  {oc}"
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Right Card
    add_card(s10, Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.2), WHITE_BG, CARD_BORDER)
    tb10_r = s10.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.4), Inches(3.8))
    tf10_r = tb10_r.text_frame
    tf10_r.word_wrap = True

    p = tf10_r.paragraphs[0]
    p.text = "LIVE DEMO FLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    demo_steps = [
        "1. Employee Login & Dashboard",
        "2. View Leave Balances",
        "3. Submit New Leave Request",
        "4. View Pending Application",
        "5. Manager Reviews Request",
        "6. Manager Approve / Reject Action",
        "7. Employee View Updated Status",
        "8. Demonstrate Postman API Validation"
    ]
    for ds in demo_steps:
        p = tf10_r.add_paragraph()
        p.space_before = Pt(4)
        p.text = ds
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK_TEXT

    # Bottom Tagline Banner
    banner = add_card(s10, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), NAVY_TITLE, NAVY_TITLE)
    tb_tag = s10.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.7))
    tf_tag = tb_tag.text_frame
    
    p = tf_tag.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CampusLoop — Simplifying Employee Leave Management Through Technology"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = tf_tag.add_paragraph()
    p.space_before = Pt(4)
    p.alignment = PP_ALIGN.CENTER
    p.text = "THANK YOU!  |  Q & A SESSION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_BOX_BG

    out_path = "CampusLoop_Presentation.pptx"
    prs.save(out_path)
    print(f"Perfect presentation saved to {out_path}")

if __name__ == "__main__":
    build_perfect_deck()
